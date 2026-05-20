import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO
from pptx.util import Emu

# Posição da imagem original no slide 15
IMG_LEFT   = Emu(828408)
IMG_TOP    = Emu(1601965)
IMG_WIDTH  = Emu(18447283)
IMG_HEIGHT = Emu(8105419)

# Cores BWAG
NAVY   = '#1C0845'
BRANCO = '#FFFFFF'
CINZA  = '#E8E8EC'
CINZA2 = '#D0D0D8'
TEXTO  = '#0A0420'
ROYAL  = '#2A1A8C'
BORDA  = '#CCCCCC'

# Proporções das colunas: Fundo | Valor (R$) | Part% | Mês | Ano | 12M | 24M
COL_X = [0.0, 0.38, 0.55, 0.63, 0.71, 0.79, 0.88, 1.0]

FIG_W = 18.0
FIG_H = FIG_W * (8105419 / 18447283)   # preserva aspect ratio exato


def _pct(v, decimais=2):
    if v is None:
        return "N/A"
    return f"{v * 100:.{decimais}f}%"


def _moeda(v):
    """Formata valor em R$ com separador de milhar."""
    if v is None:
        return "-"
    return f"R$ {v:,.2f}"


def _build_rows(fundos, perf_total, pct_cdi, titulo):
    rows = []

    # Cabeçalho título
    rows.append({"tipo": "titulo", "cols": [titulo, "", "", "", "", "", ""]})

    # Cabeçalho colunas
    rows.append({"tipo": "header",
                 "cols": ["Fundo", "Valor do Ativo (R$)", "Part. %",
                          "MÊS", "ANO", "12M", "24M"]})

    # Linhas dos fundos
    for f in fundos:
        rows.append({"tipo": "fundo",
                     "cols": [
                         f["nome"],
                         _moeda(f["financeiro"]),
                         _pct(f["pl"] / 100),
                         _pct(f["mes"]),
                         _pct(f["ano"]),
                         _pct(f["m12"]),
                         _pct(f["m24"]),
                     ]})

    # Separador
    rows.append({"tipo": "sep"})

    # Total
    total_fin = sum(f["financeiro"] for f in fundos)
    rows.append({"tipo": "total",
                 "cols": [
                     "Total",
                     _moeda(total_fin),
                     "100,00%",
                     _pct(perf_total.get("mes")),
                     _pct(perf_total.get("ano")),
                     _pct(perf_total.get("m12")),
                     _pct(perf_total.get("m24")),
                 ]})

    # % CDI
    rows.append({"tipo": "cdi",
                 "cols": [
                     "% CDI",
                     "",
                     "",
                     _pct(pct_cdi.get("mes")),
                     _pct(pct_cdi.get("ano")),
                     _pct(pct_cdi.get("m12")),
                     _pct(pct_cdi.get("m24")),
                 ]})

    return rows


def _altura(tipo):
    alturas = {
        "titulo": 1.8,
        "header": 1.6,
        "fundo":  1.1,
        "sep":    0.4,
        "total":  1.3,
        "cdi":    1.3,
    }
    return alturas.get(tipo, 1.1)


def _estilo(tipo):
    return {
        "titulo": {"bg": NAVY,   "fg": BRANCO, "bold": True,  "fs": 6.5,  "align": "center"},
        "header": {"bg": NAVY,   "fg": BRANCO, "bold": True,  "fs": 6.0,  "align": "center"},
        "fundo":  {"bg": BRANCO, "fg": TEXTO,  "bold": False, "fs": 5.5,  "align": "left"},
        "sep":    {"bg": CINZA,  "fg": TEXTO,  "bold": False, "fs": 5.0,  "align": "left"},
        "total":  {"bg": CINZA2, "fg": TEXTO,  "bold": True,  "fs": 6.0,  "align": "left"},
        "cdi":    {"bg": CINZA,  "fg": ROYAL,  "bold": False, "fs": 6.0,  "align": "left"},
    }.get(tipo, {"bg": BRANCO, "fg": TEXTO, "bold": False, "fs": 5.5, "align": "left"})


def _renderizar(rows):
    total_h = sum(_altura(r["tipo"]) for r in rows)

    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H), dpi=200)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, total_h)
    ax.axis('off')
    fig.patch.set_facecolor(BRANCO)
    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)

    y = total_h

    for row in rows:
        tipo = row["tipo"]
        est  = _estilo(tipo)
        h    = _altura(tipo)
        yb   = y - h

        cols = row.get("cols", [""] * 7)

        if tipo == "sep":
            ax.add_patch(patches.Rectangle((0, yb), 1, h, facecolor=CINZA, edgecolor='none'))
            ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.2)
            y = yb
            continue

        if tipo == "titulo":
            # Título ocupa toda a largura
            ax.add_patch(patches.Rectangle((0, yb), 1, h, facecolor=est["bg"], edgecolor=BORDA, linewidth=0.3))
            ax.text(0.5, yb + h / 2, cols[0],
                    color=est["fg"], fontsize=est["fs"],
                    ha='center', va='center', fontweight='bold')
            ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.3)
            y = yb
            continue

        # Fundo de cada célula
        for c in range(7):
            x0 = COL_X[c]
            w  = COL_X[c + 1] - x0
            ax.add_patch(patches.Rectangle((x0, yb), w, h,
                                           facecolor=est["bg"], edgecolor=BORDA, linewidth=0.2))

        # Texto
        alinhamentos = ['left', 'right', 'center', 'center', 'center', 'center', 'center']
        pad_x        = [0.008, -0.008,  0,       0,       0,       0,       0      ]

        for c in range(7):
            txt = cols[c] if c < len(cols) else ""
            if not txt:
                continue
            ha = alinhamentos[c]
            if ha == 'left':
                cx = COL_X[c] + pad_x[c]
            elif ha == 'right':
                cx = COL_X[c + 1] + pad_x[c]
            else:
                cx = (COL_X[c] + COL_X[c + 1]) / 2

            cy = yb + h / 2
            ax.text(cx, cy, txt,
                    color=est["fg"], fontsize=est["fs"],
                    ha=ha, va='center',
                    fontweight='bold' if est["bold"] else 'normal',
                    clip_on=True)

        ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.2)
        y = yb

    # Borda externa
    ax.add_patch(patches.Rectangle((0, 0), 1, total_h, fill=False,
                                   edgecolor=BORDA, linewidth=0.8))

    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf


def build_multimercado_sabia(slide, dados_sabia, config_multimercado):
    titulo = config_multimercado.get("sabia_titulo",
                                     "BWAG SABIÁ FIC MULTIMERCADO CRÉDITO PRIVADO")

    fundos     = dados_sabia["fundos"]
    perf_total = dados_sabia["perf_total"]
    pct_cdi    = dados_sabia["pct_cdi"]
    data_base  = dados_sabia["data_base"]

    rows = _build_rows(fundos, perf_total, pct_cdi, titulo)
    img  = _renderizar(rows)

    # Remove imagem original
    for shape in list(slide.shapes):
        if shape.name == "Imagem 3":
            shape._element.getparent().remove(shape._element)
            break

    # Insere nova imagem
    slide.shapes.add_picture(img, IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)

    # Atualiza rodapé com data base
    for shape in slide.shapes:
        if shape.name == "Espaço Reservado para Rodapé 4" and hasattr(shape, "text_frame"):
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text  # preserva; atualiza abaixo
                # Substitui data no texto existente
                import re
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = re.sub(r'\d{2}/\d{2}/\d{4}', data_base, run.text)
            except Exception:
                pass

    print(f"  Multimercado SABIA: {len(fundos)} fundos | data base {data_base}")
