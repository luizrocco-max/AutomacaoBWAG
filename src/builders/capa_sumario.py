from pptx import Presentation
from pptx.util import Pt
from datetime import date
import calendar
import locale

MESES_PT = {
    1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
    5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
    9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro"
}

def _titulo_periodo(competencia: str) -> str:
    """'2026-05' → 'Maio - Junho de 2026'"""
    ano, mes = int(competencia[:4]), int(competencia[5:7])
    mes_seguinte = mes % 12 + 1
    ano_seguinte = ano + 1 if mes == 12 else ano
    return f"{MESES_PT[mes]} - {MESES_PT[mes_seguinte]} de {ano_seguinte}"

def _atualizar_texto(shape, novo_texto: str):
    """Atualiza o texto de um shape preservando a formatação do primeiro run."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = novo_texto
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.text = novo_texto

def build_capa(slide, competencia: str):
    """Atualiza o título da capa com o período correto."""
    titulo = _titulo_periodo(competencia)
    for shape in slide.shapes:
        if shape.name == "Título 2" and shape.has_text_frame:
            _atualizar_texto(shape, titulo)
            print(f"  Capa: titulo atualizado: '{titulo}'")
            return
    print("  Capa: shape 'Título 2' não encontrado")

def build_sumario(slide, secoes: dict):
    """
    Atualiza os números de slide no sumário.

    secoes: dict mapeando nome da seção → número do slide inicial
    Ex: {"Overview": 3, "Performance": 5, ...}
    """
    rectangles = [s for s in slide.shapes if s.name == "Rectangle 19" and s.has_text_frame]
    numeros = list(secoes.values())

    for i, rect in enumerate(rectangles):
        if i < len(numeros):
            _atualizar_texto(rect, str(numeros[i]))

    print(f"  Sumário: {len(rectangles)} números de slide atualizados")
