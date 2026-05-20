import yaml
import os
from dotenv import load_dotenv
from pptx import Presentation

from src.builders.capa_sumario import build_capa, build_sumario
from src.builders.overview import build_overview
from src.builders.performance import build_performance
from src.builders.multimercado_sabia import build_multimercado_sabia
from src.builders.renda_variavel_maitaca import build_renda_variavel_maitaca, ler_mapeamento_nomes
from src.readers.quantum_reader import ler_quantum
from src.readers.quantum_performance_reader import ler_performance
from src.readers.btg_pdf_reader import carregar_overrides_btg
from src.readers.btg_sabia_reader import ler_sabia, encontrar_pdf_sabia
from src.readers.btg_carteira_reader import ler_carteira_btg, encontrar_pdf_fundo
from src.readers.quantum_beta_reader import ler_quantum_beta

load_dotenv()

def carregar_config(caminho: str = "config/config.yaml") -> dict:
    with open(caminho, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)

def main():
    config = carregar_config()
    competencia = config["competencia"]
    print(f"\nGerando apresentacao - competencia: {competencia}\n")

    prs = Presentation(config["caminhos"]["template"])

    print("[ Capa ]")
    build_capa(prs.slides[0], competencia)

    print("[ Sumario ]")
    build_sumario(prs.slides[1], config["sumario"])

    print("[ Overview ]")
    cfg_ov = config["overview"]
    dados_quantum = ler_quantum(
        caminho=cfg_ov["arquivo_quantum"],
        mapeamento_nomes=cfg_ov.get("nomes_curtos", {}),
    )
    build_overview(
        slide=prs.slides[3],
        dados_quantum=dados_quantum,
        valor_fundos=cfg_ov["valor_fundos"],
        valor_direto=cfg_ov["valor_direto"],
    )

    print("[ Performance ]")
    cfg_perf = config["performance"]
    dados_perf = ler_performance(cfg_perf["arquivo_quantum"])
    overrides_btg = carregar_overrides_btg(config["caminhos"]["input_pdf"])
    build_performance(
        slide=prs.slides[5],
        dados_quantum=dados_perf,
        overrides_btg=overrides_btg,
        config_performance=cfg_perf,
    )

    print("[ Multimercado - SABIA ]")
    cfg_mult = config.get("multimercado", {})
    pdf_sabia = encontrar_pdf_sabia(config["caminhos"]["input_pdf"])
    if pdf_sabia:
        dados_sabia = ler_sabia(pdf_sabia)
        build_multimercado_sabia(
            slide=prs.slides[14],
            dados_sabia=dados_sabia,
            config_multimercado=cfg_mult,
        )
    else:
        print("  Aviso: PDF do SABIA nao encontrado em input_pdf/")

    print("[ Renda Variável - MAITACA ]")
    cfg_rv = config.get("renda_variavel", {})
    pdf_maitaca = encontrar_pdf_fundo(config["caminhos"]["input_pdf"], "maitaca")
    if pdf_maitaca and cfg_rv.get("arquivo_quantum_beta"):
        dados_maitaca  = ler_carteira_btg(pdf_maitaca)
        dados_beta     = ler_quantum_beta(cfg_rv["arquivo_quantum_beta"])
        mapa_nomes     = ler_mapeamento_nomes(
            os.path.join(config["caminhos"]["input_excel"], "mapeamento_nomes_maitaca.xlsx")
        )
        build_renda_variavel_maitaca(
            slide=prs.slides[21],
            dados_pdf=dados_maitaca,
            dados_beta=dados_beta,
            config_maitaca={
                "titulo":            cfg_rv.get("maitaca_titulo", "MAITACA AÇÕES FIC AÇÕES"),
                "nomes_pdf_quantum": mapa_nomes,
            },
        )
    else:
        print("  Aviso: PDF ou Quantum Beta do MAITACA não encontrado")

    pasta_saida = os.path.join(config["caminhos"]["output"], competencia)
    os.makedirs(pasta_saida, exist_ok=True)
    arquivo_saida = os.path.join(pasta_saida, f"comite_{competencia}.pptx")
    prs.save(arquivo_saida)
    print(f"\nApresentacao salva em: {arquivo_saida}")

if __name__ == "__main__":
    main()
