import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO
from pptx.util import Emu

# Posição da imagem original no slide
IMG_LEFT   = Emu(1916248)
IMG_TOP    = Emu(1352780)
IMG_WIDTH  = Emu(16271603)
IMG_HEIGHT = Emu(8603790)

# Cores BWAG (matplotlib)
NAVY    = '#1C0845'
ROYAL   = '#2A1A8C'
BRANCO  = '#FFFFFF'
CINZA   = '#E8E8EC'
TEXTO   = '#0A0420'
LARANJA = '#C05010'
BORDA   = '#CCCCCC'

# Proporções das colunas (0 a 1)
COL_X = [0.0, 0.515, 0.636, 0.757, 0.878, 1.0]

def _pct(v):
    if v is None:
        return "N/A"
    return f"{v * 100:.2f}%"

def _calcular_bench(fundo_vals, bench_vals, tipo):
    res = {}
    for p in ["mes", "ano", "m12", "m24"]:
        f = fundo_vals.get(p)
        b = bench_vals.get(p) if bench_vals else None
        if f is None or b is None or b == 0:
            res[p] = None
        elif tipo == "cdi":
            res[p] = f / b
        else:
            res[p] = f - b
    return res

def _aplicar_btg(fundo_vals, overrides_btg, nome_fundo):
    for nome_btg, dados in overrides_btg.items():
        if nome_btg.upper() in nome_fundo.upper() or nome_fundo.upper() in nome_btg.upper():
            print(f"    Override BTG aplicado: {nome_btg}")
            return {k: dados[k] for k in ["mes", "ano", "m12", "m24"]}
    return fundo_vals

def _preparar_secao(lista, q_fundos, overrides_btg, nomes_curtos, bench_vals, tipo):
    linhas = []
    for nome in lista:
        if nome not in q_fundos:
            continue
        vals  = _aplicar_btg(q_fundos[nome], overrides_btg, nome)
        bench = _calcular_bench(vals, bench_vals, tipo)
        curto = nomes_curtos.get(nome, nome)
        linhas.append((curto, vals, bench))
    return linhas

def _construir_rows(sec_cdi, sec_ibov, benchmarks):
    """Monta lista de rows com dados e estilo."""
    rows = []

    # Header 1
    rows.append({"tipo": "h1",
                 "cols": ["", "RENTABILIDADE (%)", "", "", ""],
                 "merge": True})
    # Header 2
    rows.append({"tipo": "h2",
                 "cols": ["FUNDO", "MÊS", "ANO", "12M", "24M"]})

    # CDI
    for nome, vals, bench in sec_cdi:
        rows.append({"tipo": "fundo",
                     "cols": [nome,
                               _pct(vals.get("mes")), _pct(vals.get("ano")),
                               _pct(vals.get("m12")), _pct(vals.get("m24"))]})
        rows.append({"tipo": "cdi",
                     "cols": ["% CDI",
                               _pct(bench.get("mes")), _pct(bench.get("ano")),
                               _pct(bench.get("m12")), _pct(bench.get("m24"))]})

    rows.append({"tipo": "sep"})

    # Ibov
    for nome, vals, bench in sec_ibov:
        rows.append({"tipo": "fundo",
                     "cols": [nome,
                               _pct(vals.get("mes")), _pct(vals.get("ano")),
                               _pct(vals.get("m12")), _pct(vals.get("m24"))]})
        rows.append({"tipo": "ibov",
                     "cols": ["α Ibov",
                               _pct(bench.get("mes")), _pct(bench.get("ano")),
                               _pct(bench.get("m12")), _pct(bench.get("m24"))]})

    rows.append({"tipo": "sep"})

    # Benchmarks
    bench_nomes = ["CDI", "Dólar", "Ibovespa", "IPCA"]
    for nome in bench_nomes:
        bv = benchmarks.get(nome, {})
        rows.append({"tipo": "bench",
                     "cols": [nome,
                               _pct(bv.get("mes")), _pct(bv.get("ano")),
                               _pct(bv.get("m12")), _pct(bv.get("m24"))]})
    return rows

def _altura_row(tipo):
    alturas = {"h1": 2.0, "h2": 1.8, "fundo": 1.4, "cdi": 1.4,
               "ibov": 1.4, "sep": 0.6, "bench": 1.4}
    return alturas.get(tipo, 1.4)

def _estilo(tipo):
    estilos = {
        "h1":    {"bg": NAVY,   "fg": BRANCO, "bold": True,  "italic": False, "fs": 7.5},
        "h2":    {"bg": NAVY,   "fg": BRANCO, "bold": True,  "italic": False, "fs": 7.0},
        "fundo": {"bg": BRANCO, "fg": TEXTO,  "bold": False, "italic": False, "fs": 6.5},
        "cdi":   {"bg": BRANCO, "fg": ROYAL,  "bold": False, "italic": True,  "fs": 6.5},
        "ibov":  {"bg": BRANCO, "fg": LARANJA,"bold": False, "italic": True,  "fs": 6.5},
        "sep":   {"bg": CINZA,  "fg": TEXTO,  "bold": False, "italic": False, "fs": 6.0},
        "bench": {"bg": CINZA,  "fg": TEXTO,  "bold": False, "italic": True,  "fs": 6.5},
    }
    return estilos.get(tipo, estilos["fundo"])

def _renderizar_imagem(rows):
    """Renderiza a tabela como PNG e retorna BytesIO."""
    # Aspect ratio exato da imagem original no slide
    FIG_W = 16.0
    FIG_H = FIG_W * (8603790 / 16271603)  # = 8.463 inches

    total_h = sum(_altura_row(r["tipo"]) for r in rows)

    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H), dpi=200)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, total_h)
    ax.axis('off')
    fig.patch.set_facecolor(BRANCO)
    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)

    y = total_h  # começa do topo

    for row in rows:
        tipo  = row["tipo"]
        est   = _estilo(tipo)
        h     = _altura_row(tipo)
        y_bot = y - h

        if tipo == "sep":
            ax.add_patch(patches.Rectangle(
                (0, y_bot), 1, h, facecolor=CINZA, edgecolor='none'))
            ax.plot([0, 1], [y_bot, y_bot], color=BORDA, linewidth=0.3)
            y = y_bot
            continue

        cols = row.get("cols", [""] * 5)

        # Fundo das células
        for c in range(5):
            x0 = COL_X[c]
            w  = COL_X[c+1] - COL_X[c]
            ax.add_patch(patches.Rectangle(
                (x0, y_bot), w, h,
                facecolor=est["bg"], edgecolor=BORDA, linewidth=0.3))

        # Texto nas células
        alinhamentos = ['left', 'center', 'center', 'center', 'center']
        offsetx      = [0.008,  0,         0,         0,         0      ]

        if tipo == "h1":
            # Col 0 vazia, cols 1-4 mescla (texto centralizado)
            cx = (COL_X[1] + COL_X[5]) / 2
            cy = y_bot + h / 2
            ax.text(cx, cy, cols[1], color=est["fg"],
                    fontsize=est["fs"], ha='center', va='center',
                    fontweight='bold', style='normal')
        else:
            for c in range(5):
                cx = COL_X[c] + offsetx[c] + (COL_X[c+1] - COL_X[c]) * (0.02 if c == 0 else 0.5)
                cy = y_bot + h / 2
                ax.text(cx, cy, cols[c],
                        color=est["fg"], fontsize=est["fs"],
                        ha=alinhamentos[c], va='center',
                        fontweight='bold' if est["bold"] else 'normal',
                        style='italic' if est["italic"] else 'normal',
                        clip_on=True)

        # Linha separadora horizontal
        ax.plot([0, 1], [y_bot, y_bot], color=BORDA, linewidth=0.3)

        y = y_bot

    # Borda externa
    ax.add_patch(patches.Rectangle(
        (0, 0), 1, total_h, fill=False, edgecolor=BORDA, linewidth=0.8))

    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf

def build_performance(slide, dados_quantum, overrides_btg, config_performance):
    fundos_cdi   = config_performance["fundos_cdi"]
    fundos_ibov  = config_performance["fundos_ibov"]
    nomes_curtos = config_performance.get("nomes_curtos", {})

    q_fundos     = dados_quantum["fundos"]
    q_benchmarks = dados_quantum["benchmarks"]
    cdi_vals     = q_benchmarks.get("CDI", {})
    ibov_vals    = q_benchmarks.get("Ibovespa", {})

    sec_cdi  = _preparar_secao(fundos_cdi,  q_fundos, overrides_btg, nomes_curtos, cdi_vals,  "cdi")
    sec_ibov = _preparar_secao(fundos_ibov, q_fundos, overrides_btg, nomes_curtos, ibov_vals, "ibov")

    rows = _construir_rows(sec_cdi, sec_ibov, q_benchmarks)
    img  = _renderizar_imagem(rows)

    # Remove imagem original
    for shape in list(slide.shapes):
        if shape.name == "Imagem 4":
            shape._element.getparent().remove(shape._element)
            break

    # Insere nova imagem na mesma posição
    slide.shapes.add_picture(img, IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)

    print(f"  Performance: {len(sec_cdi)} fundos CDI + {len(sec_ibov)} fundos Ibov [imagem]")
