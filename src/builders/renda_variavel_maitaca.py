import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO
from pptx.util import Emu
import re

# Posição da imagem no slide 22
IMG_LEFT   = Emu(828408)
IMG_TOP    = Emu(1601965)
IMG_WIDTH  = Emu(18447283)
IMG_HEIGHT = Emu(8105419)

NAVY   = '#1C0845'
BRANCO = '#FFFFFF'
CINZA  = '#E8E8EC'
CINZA2 = '#D0D0D8'
TEXTO  = '#0A0420'
ROYAL  = '#2A1A8C'
LARANJA = '#C05010'
BORDA  = '#CCCCCC'

# Colunas: Fundo | Valor (R$) | Part% | Mês | Ano | 12M | 24M | Beta(6m) | Beta(12m)
COL_X = [0.0, 0.32, 0.46, 0.53, 0.61, 0.69, 0.77, 0.855, 0.928, 1.0]

FIG_W = 18.0
FIG_H = FIG_W * (8105419 / 18447283)


def _pct(v, decimais=2):
    if v is None:
        return "N/D"
    return f"{v * 100:.{decimais}f}%"


def _beta(v):
    if v is None:
        return "N/D"
    return f"{v:.2f}"


def _moeda(v):
    if v is None:
        return "-"
    return f"R$ {v:,.2f}"


def ler_mapeamento_nomes(caminho: str) -> dict:
    """Lê Excel com colunas Nome_PDF e Nome_Completo."""
    import pandas as pd, os
    if not os.path.exists(caminho):
        return {}
    df = pd.read_excel(caminho)
    mapa = {}
    for _, row in df.iterrows():
        pdf  = str(row.get("Nome_PDF", "")).strip()
        comp = str(row.get("Nome_Completo", "")).strip()
        if pdf and comp and comp.lower() != "nan":
            mapa[pdf] = comp
    return mapa


def _preparar_fundos(fundos_pdf, dados_beta, mapeamento):
    """
    Junta carteira PDF com Beta Quantum.
    Agrupa fundos com mesmo nome Quantum (ex: dois CAPSIGMA).
    """
    quantum_nomes = list(dados_beta["fundos"].keys())
    grupos = {}  # nome_display → dados acumulados

    for f in fundos_pdf:
        nome_pdf = f["nome"]

        # Resolve nome display e chave Quantum
        nome_q = mapeamento.get(nome_pdf)
        if nome_q and nome_q in dados_beta["fundos"]:
            nome_display = nome_q
            beta_dados = dados_beta["fundos"][nome_q]
        elif nome_q:
            nome_display = nome_q  # display name sem beta
            beta_dados = {}
        else:
            nome_display = nome_pdf
            beta_dados = {}

        if nome_display not in grupos:
            grupos[nome_display] = {
                "pl":         f["pl"],
                "financeiro": f["financeiro"],
                "mes":        f["mes"],
                "ano":        f["ano"],
                "m12":        f["m12"],
                "m24":        f["m24"],
                "beta6m":     beta_dados.get("beta6m"),
                "beta12m":    beta_dados.get("beta12m"),
            }
        else:
            # Agrega (ex: dois CAPSIGMA)
            g = grupos[nome_display]
            g["pl"]         += f["pl"]
            g["financeiro"] += f["financeiro"]
            # Média ponderada das rentabilidades
            # (simplificação: usa a rentabilidade do maior)

    resultado = [{"nome": k, **v} for k, v in grupos.items()]
    resultado.sort(key=lambda x: x["pl"], reverse=True)
    return resultado


def _build_rows(fundos, perf_total, ibovespa, titulo):
    rows = []

    rows.append({"tipo": "titulo",
                 "cols": [titulo, "", "", "", "", "", "", ""]})
    rows.append({"tipo": "header",
                 "cols": ["Fundo", "Valor do Ativo (R$)", "Part. %",
                          "MÊS", "ANO", "12M", "24M", "Beta\n(6m)", "Beta\n(12m)"]})

    for f in fundos:
        rows.append({"tipo": "fundo",
                     "cols": [
                         f["nome"],
                         _moeda(f["financeiro"]),
                         _pct(f["pl"] / 100),
                         _pct(f["mes"]),
                         _pct(f["ano"]),
                         _pct(f["m12"]),
                         _pct(f["m24"]),
                         _beta(f["beta6m"]),
                         _beta(f["beta12m"]),
                     ]})

    rows.append({"tipo": "sep"})

    # Total
    total_fin = sum(f["financeiro"] for f in fundos)
    total_b6  = _beta(next((f["beta6m"]  for f in fundos if f["nome"] == titulo.split()[0]), None))

    # Média ponderada de Beta para o total (fund-level beta)
    # Usa os valores de Beta do próprio fundo MAITACA no Quantum (se disponível)
    maitaca_q = next((f for f in fundos if "MAITACA" in f["nome"].upper()), None)

    rows.append({"tipo": "total",
                 "cols": [
                     "Total",
                     _moeda(total_fin),
                     "100,00%",
                     _pct(perf_total.get("mes")),
                     _pct(perf_total.get("ano")),
                     _pct(perf_total.get("m12")),
                     _pct(perf_total.get("m24")),
                     "", "",
                 ]})

    # α Ibov = fundo - ibovespa
    def _alpha(key):
        f = perf_total.get(key)
        b = ibovespa.get(key)
        if f is None or b is None:
            return "N/D"
        return _pct(f - b)

    rows.append({"tipo": "ibov",
                 "cols": [
                     "α Ibov",
                     "", "",
                     _alpha("mes"),
                     _alpha("ano"),
                     _alpha("m12"),
                     _alpha("m24"),
                     "", "",
                 ]})

    return rows


def _altura(tipo):
    return {"titulo": 1.8, "header": 1.8, "fundo": 1.05,
            "sep": 0.4, "total": 1.3, "ibov": 1.3}.get(tipo, 1.05)


def _estilo(tipo):
    return {
        "titulo": {"bg": NAVY,   "fg": BRANCO, "bold": True,  "fs": 6.5},
        "header": {"bg": NAVY,   "fg": BRANCO, "bold": True,  "fs": 5.8},
        "fundo":  {"bg": BRANCO, "fg": TEXTO,  "bold": False, "fs": 5.2},
        "sep":    {"bg": CINZA,  "fg": TEXTO,  "bold": False, "fs": 5.0},
        "total":  {"bg": CINZA2, "fg": TEXTO,  "bold": True,  "fs": 5.8},
        "ibov":   {"bg": CINZA,  "fg": LARANJA,"bold": False, "fs": 5.8},
    }.get(tipo, {"bg": BRANCO, "fg": TEXTO, "bold": False, "fs": 5.2})


def _renderizar(rows):
    total_h = sum(_altura(r["tipo"]) for r in rows)
    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H), dpi=200)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, total_h)
    ax.axis('off')
    fig.patch.set_facecolor(BRANCO)
    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)

    y = total_h
    n_cols = len(COL_X) - 1

    for row in rows:
        tipo = row["tipo"]
        est  = _estilo(tipo)
        h    = _altura(tipo)
        yb   = y - h
        cols = row.get("cols", [""] * n_cols)

        if tipo == "sep":
            ax.add_patch(patches.Rectangle((0, yb), 1, h,
                                           facecolor=CINZA, edgecolor='none'))
            ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.2)
            y = yb
            continue

        if tipo == "titulo":
            ax.add_patch(patches.Rectangle((0, yb), 1, h,
                                           facecolor=est["bg"], edgecolor=BORDA, linewidth=0.3))
            ax.text(0.5, yb + h / 2, cols[0],
                    color=est["fg"], fontsize=est["fs"],
                    ha='center', va='center', fontweight='bold')
            ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.2)
            y = yb
            continue

        # Células
        for c in range(n_cols):
            x0 = COL_X[c]
            w  = COL_X[c + 1] - x0
            ax.add_patch(patches.Rectangle((x0, yb), w, h,
                                           facecolor=est["bg"], edgecolor=BORDA, linewidth=0.2))

        aligns = ['left', 'right', 'center', 'center', 'center',
                  'center', 'center', 'center', 'center']
        pads   = [0.007, -0.006, 0, 0, 0, 0, 0, 0, 0]

        for c in range(min(n_cols, len(cols))):
            txt = cols[c]
            if not txt:
                continue
            ha = aligns[c]
            if ha == 'left':
                cx = COL_X[c] + pads[c]
            elif ha == 'right':
                cx = COL_X[c + 1] + pads[c]
            else:
                cx = (COL_X[c] + COL_X[c + 1]) / 2

            ax.text(cx, yb + h / 2, txt,
                    color=est["fg"], fontsize=est["fs"],
                    ha=ha, va='center',
                    fontweight='bold' if est["bold"] else 'normal',
                    clip_on=True)

        ax.plot([0, 1], [yb, yb], color=BORDA, linewidth=0.2)
        y = yb

    ax.add_patch(patches.Rectangle((0, 0), 1, total_h, fill=False,
                                   edgecolor=BORDA, linewidth=0.8))

    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf


def build_renda_variavel_maitaca(slide, dados_pdf, dados_beta, config_maitaca):
    titulo    = config_maitaca.get("titulo", "MAITACA AÇÕES FIC AÇÕES")
    mapeamento = config_maitaca.get("nomes_pdf_quantum", {})

    fundos     = _preparar_fundos(dados_pdf["fundos"], dados_beta, mapeamento)
    perf_total = dados_pdf["perf_total"]
    ibovespa   = dados_beta.get("ibovespa", {})
    data_base  = dados_pdf["data_base"]

    rows = _build_rows(fundos, perf_total, ibovespa, titulo)
    img  = _renderizar(rows)

    # Descobre o nome da imagem no slide
    img_nome = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            img_nome = shape.name
            break

    for shape in list(slide.shapes):
        if shape.name == img_nome:
            shape._element.getparent().remove(shape._element)
            break

    slide.shapes.add_picture(img, IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)

    # Atualiza rodapé
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "data base" in shape.text.lower():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = re.sub(r'\d{2}/\d{2}/\d{4}', data_base, run.text)

    print(f"  Maitaca: {len(fundos)} fundos | data base {data_base}")
